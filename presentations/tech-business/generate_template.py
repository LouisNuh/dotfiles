#!/usr/bin/env python3
"""
Generate a tech-business style PowerPoint template.
This script creates a PPTX file with professional styling for business presentations.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme (Tech-Business style)
COLOR_DEEP_TECH_BLUE = RGBColor(11, 59, 113)     # #0B3B71
COLOR_CYAN_GREEN = RGBColor(12, 170, 170)        # #0CAAAA
COLOR_BG_GRAY = RGBColor(245, 247, 250)          # #F5F7FA
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_DARK_TEXT = RGBColor(51, 51, 51)
COLOR_LIGHT_GRAY = RGBColor(200, 200, 200)

# Dimensions (16:9 aspect ratio)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

# Margins (2.54cm = 1 inch left/right, 1.9cm = 0.748 inch top/bottom)
MARGIN_LEFT = Inches(1.0)
MARGIN_RIGHT = Inches(1.0)
MARGIN_TOP = Inches(0.748)
MARGIN_BOTTOM = Inches(0.748)

def create_presentation():
    """Create a new PowerPoint presentation with 16:9 aspect ratio."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs

def add_title_slide(prs):
    """Add a title slide with tech-business styling."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG_GRAY
    background.line.color.rgb = COLOR_BG_GRAY
    
    # Left accent bar
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        Inches(0.2), SLIDE_HEIGHT
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_CYAN_GREEN
    accent_bar.line.color.rgb = COLOR_CYAN_GREEN
    
    # Title
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(0.3),
        Inches(1.5),
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT - Inches(0.3),
        Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.text = "科技商业演示模板"
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_DEEP_TECH_BLUE
    title_para.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(0.3),
        Inches(3.2),
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT - Inches(0.3),
        Inches(1.0)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "Tech-Business Presentation Template"
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = COLOR_DARK_TEXT
    subtitle_para.alignment = PP_ALIGN.LEFT
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_WHITE
    background.line.color.rgb = COLOR_WHITE
    
    # Top accent bar
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        SLIDE_WIDTH, Inches(0.15)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_CYAN_GREEN
    accent_bar.line.color.rgb = COLOR_CYAN_GREEN
    
    # Title
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT,
        MARGIN_TOP + Inches(0.2),
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
        Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_DEEP_TECH_BLUE
    title_para.alignment = PP_ALIGN.LEFT
    
    # Content area with bullet points
    content_top = MARGIN_TOP + Inches(1.2)
    content_box = slide.shapes.add_textbox(
        MARGIN_LEFT + Inches(0.2),
        content_top,
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT - Inches(0.2),
        SLIDE_HEIGHT - content_top - MARGIN_BOTTOM
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            content_frame.add_paragraph()
        para = content_frame.paragraphs[i]
        para.text = item
        para.font.size = Pt(24)
        para.font.color.rgb = COLOR_DARK_TEXT
        para.level = 0
        para.space_before = Pt(10)
        para.space_after = Pt(10)
        
        # Add bullet point (circular)
        para.text = "• " + item
    
    return slide

def add_two_column_slide(prs, title, left_items, right_items):
    """Add a two-column content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_WHITE
    background.line.color.rgb = COLOR_WHITE
    
    # Top accent bar
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        SLIDE_WIDTH, Inches(0.15)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_CYAN_GREEN
    accent_bar.line.color.rgb = COLOR_CYAN_GREEN
    
    # Title
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT,
        MARGIN_TOP + Inches(0.2),
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
        Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_DEEP_TECH_BLUE
    
    # Content area
    content_top = MARGIN_TOP + Inches(1.2)
    content_width = (SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT - Inches(0.3)) / 2
    
    # Left column
    left_box = slide.shapes.add_textbox(
        MARGIN_LEFT,
        content_top,
        content_width,
        SLIDE_HEIGHT - content_top - MARGIN_BOTTOM
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, item in enumerate(left_items):
        if i > 0:
            left_frame.add_paragraph()
        para = left_frame.paragraphs[i]
        para.text = "• " + item
        para.font.size = Pt(20)
        para.font.color.rgb = COLOR_DARK_TEXT
        para.space_before = Pt(8)
        para.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(
        MARGIN_LEFT + content_width + Inches(0.3),
        content_top,
        content_width,
        SLIDE_HEIGHT - content_top - MARGIN_BOTTOM
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, item in enumerate(right_items):
        if i > 0:
            right_frame.add_paragraph()
        para = right_frame.paragraphs[i]
        para.text = "• " + item
        para.font.size = Pt(20)
        para.font.color.rgb = COLOR_DARK_TEXT
        para.space_before = Pt(8)
        para.space_after = Pt(8)
    
    return slide

def add_conclusion_slide(prs, title, conclusion_text):
    """Add a conclusion/summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background with gradient effect (simulated with two colors)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_DEEP_TECH_BLUE
    background.line.color.rgb = COLOR_DEEP_TECH_BLUE
    
    # Accent shape
    accent_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        SLIDE_WIDTH - Inches(3.5),
        Inches(1.0),
        Inches(3.0),
        Inches(3.5)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = COLOR_CYAN_GREEN
    accent_shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT,
        Inches(1.5),
        SLIDE_WIDTH - Inches(4.0),
        Inches(1.2)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(56)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_WHITE
    
    # Conclusion text
    text_box = slide.shapes.add_textbox(
        MARGIN_LEFT,
        Inches(3.0),
        SLIDE_WIDTH - Inches(4.0),
        Inches(2.0)
    )
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_para = text_frame.paragraphs[0]
    text_para.text = conclusion_text
    text_para.font.size = Pt(24)
    text_para.font.color.rgb = COLOR_WHITE
    
    return slide

def add_agenda_slide(prs):
    """Add an agenda/table of contents slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_WHITE
    background.line.color.rgb = COLOR_WHITE
    
    # Top accent bar
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        SLIDE_WIDTH, Inches(0.15)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_CYAN_GREEN
    accent_bar.line.color.rgb = COLOR_CYAN_GREEN
    
    # Title
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT,
        MARGIN_TOP + Inches(0.2),
        SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
        Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "目录"
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_DEEP_TECH_BLUE
    
    # Agenda items
    agenda_items = [
        "项目概述与背景",
        "核心功能与特性",
        "技术架构方案",
        "实施计划与时间表",
        "总结与展望"
    ]
    
    content_top = MARGIN_TOP + Inches(1.3)
    item_height = Inches(0.6)
    
    for i, item in enumerate(agenda_items):
        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            MARGIN_LEFT,
            content_top + i * item_height,
            Inches(0.4),
            Inches(0.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_DEEP_TECH_BLUE
        circle.line.color.rgb = COLOR_DEEP_TECH_BLUE
        
        # Number text
        number_frame = circle.text_frame
        number_para = number_frame.paragraphs[0]
        number_para.text = str(i + 1)
        number_para.font.size = Pt(18)
        number_para.font.bold = True
        number_para.font.color.rgb = COLOR_WHITE
        number_para.alignment = PP_ALIGN.CENTER
        number_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Item text
        item_box = slide.shapes.add_textbox(
            MARGIN_LEFT + Inches(0.6),
            content_top + i * item_height,
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT - Inches(0.6),
            item_height
        )
        item_frame = item_box.text_frame
        item_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        item_para = item_frame.paragraphs[0]
        item_para.text = item
        item_para.font.size = Pt(24)
        item_para.font.color.rgb = COLOR_DARK_TEXT
    
    return slide

def main():
    """Generate the tech-business PowerPoint template."""
    print("Creating tech-business PowerPoint template...")
    
    prs = create_presentation()
    
    # Slide 1: Title slide
    print("Adding title slide...")
    add_title_slide(prs)
    
    # Slide 2: Agenda/Table of contents
    print("Adding agenda slide...")
    add_agenda_slide(prs)
    
    # Slide 3: Content slide example
    print("Adding content slide...")
    add_content_slide(
        prs,
        "项目概述",
        [
            "项目背景：响应市场需求，提供创新解决方案",
            "目标定位：成为行业领先的技术服务平台",
            "核心价值：提升效率、降低成本、优化体验",
            "预期成果：实现业务增长30%以上"
        ]
    )
    
    # Slide 4: Two-column slide
    print("Adding two-column slide...")
    add_two_column_slide(
        prs,
        "核心功能",
        [
            "智能数据分析",
            "实时监控预警",
            "自动化流程管理",
            "多维度报表生成"
        ],
        [
            "云端协同办公",
            "移动端随时访问",
            "安全加密传输",
            "7×24小时技术支持"
        ]
    )
    
    # Slide 5: Another content slide
    print("Adding technical architecture slide...")
    add_content_slide(
        prs,
        "技术架构",
        [
            "前端：React + TypeScript 响应式设计",
            "后端：微服务架构，支持高并发处理",
            "数据库：分布式存储，保障数据安全",
            "部署：容器化部署，支持弹性扩展"
        ]
    )
    
    # Slide 6: Conclusion slide
    print("Adding conclusion slide...")
    add_conclusion_slide(
        prs,
        "谢谢！",
        "期待与您的合作\nThank you for your attention"
    )
    
    # Save the presentation
    output_path = "tech-business-template.pptx"
    prs.save(output_path)
    print(f"✓ Template saved as: {output_path}")
    
    return output_path

if __name__ == "__main__":
    main()
