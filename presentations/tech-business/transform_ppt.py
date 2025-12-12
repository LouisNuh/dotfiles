#!/usr/bin/env python3
"""
Transform existing PowerPoint to tech-business style template.
This script modifies the presentation while preserving all text as editable.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Tech-Business Color Scheme
PRIMARY_COLOR = RGBColor(11, 59, 113)      # #0B3B71
SECONDARY_COLOR = RGBColor(12, 170, 170)   # #0CAAAA
BACKGROUND_COLOR = RGBColor(245, 247, 250) # #F5F7FA
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(51, 51, 51)

def set_font_properties(text_frame, font_name="Arial", font_size=22, bold=False, color=DARK_TEXT):
    """Apply consistent font properties to text frame."""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color

def apply_tech_business_style(input_path, output_path):
    """Transform PPT to tech-business style."""
    
    print(f"Loading presentation from {input_path}...")
    prs = Presentation(input_path)
    
    # Set 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    print(f"Processing {len(prs.slides)} slide(s)...")
    
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"  Processing slide {slide_num}...")
        
        # Apply background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BACKGROUND_COLOR
        
        for shape in slide.shapes:
            # Apply styles based on shape type
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame
                
                # Detect if this is a title (longer text or first shape)
                is_title = False
                if shape.text and len(shape.text) > 10:
                    # Check if it contains keywords that suggest it's a title
                    if any(keyword in shape.text for keyword in ["挑战", "方案", "强磁场", "科学", "工程"]):
                        is_title = True
                
                # Apply font styles
                if is_title:
                    set_font_properties(text_frame, font_name="Arial", font_size=48, bold=True, color=PRIMARY_COLOR)
                else:
                    # Check if it's a bullet point or body text
                    has_bullets = False
                    for paragraph in text_frame.paragraphs:
                        if paragraph.level > 0 or (paragraph.text and paragraph.text.strip().startswith(('•', '-', '·'))):
                            has_bullets = True
                            break
                    
                    if has_bullets:
                        set_font_properties(text_frame, font_name="Arial", font_size=22, bold=False, color=DARK_TEXT)
                    else:
                        set_font_properties(text_frame, font_name="Arial", font_size=24, bold=False, color=DARK_TEXT)
                
                # Set text alignment and margins
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.margin_top = Inches(0.05)
                text_frame.margin_bottom = Inches(0.05)
                
                # Adjust line spacing
                for paragraph in text_frame.paragraphs:
                    paragraph.line_spacing = 1.2
                    paragraph.space_before = Pt(6)
                    paragraph.space_after = Pt(6)
            
            # Apply shape fill colors
            if hasattr(shape, "fill") and shape.shape_type != 14:  # Not a placeholder
                try:
                    # Use subtle colors for shapes
                    fill = shape.fill
                    if fill.type == 1:  # Solid fill
                        # Keep existing colors or apply subtle background
                        pass
                except:
                    pass
    
    print(f"Saving transformed presentation to {output_path}...")
    prs.save(output_path)
    print("Transformation complete!")

def create_master_template(output_path):
    """Create a comprehensive tech-business template with master layouts."""
    
    print(f"Creating master template at {output_path}...")
    prs = Presentation()
    
    # Set 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    print("  Creating title slide layout...")
    slide_layout = prs.slide_layouts[0]  # Title layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Apply background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND_COLOR
    
    # Title
    title = slide.shapes.title
    title.text = "科技商业演示模板"
    title.text_frame.paragraphs[0].font.name = "Arial"
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = "Tech-Business Presentation Template"
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.font.name = "Arial"
            paragraph.font.size = Pt(32)
            paragraph.font.color.rgb = SECONDARY_COLOR
            paragraph.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Content Slide with bullets
    print("  Creating content slide layout...")
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND_COLOR
    
    title = slide.shapes.title
    title.text = "主要特点 / Key Features"
    title.text_frame.paragraphs[0].font.name = "Arial"
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Content area
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    
    # Add bullet points
    items = [
        "主题色：主色 #0B3B71，辅色 #0CAAAA",
        "背景色：中性 #F5F7FA",
        "字体：中文 Noto Sans CJK / PingFang SC，英文 Montserrat / Arial",
        "标题：56-64pt 加粗",
        "正文：22-28pt",
        "版式：16:9 宽屏格式"
    ]
    
    for item in items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_TEXT
        p.line_spacing = 1.3
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    # Slide 3: Two Column Layout
    print("  Creating two-column slide layout...")
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND_COLOR
    
    title = slide.shapes.title
    title.text = "双栏布局示例 / Two-Column Layout"
    title.text_frame.paragraphs[0].font.name = "Arial"
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.0), Inches(6.0), Inches(4.5)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    p = left_frame.paragraphs[0]
    p.text = "左栏内容 / Left Column"
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    p = left_frame.add_paragraph()
    p.text = "• 要点一\n• 要点二\n• 要点三"
    p.font.name = "Arial"
    p.font.size = Pt(22)
    p.font.color.rgb = DARK_TEXT
    
    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(7.0), Inches(2.0), Inches(6.0), Inches(4.5)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    p = right_frame.paragraphs[0]
    p.text = "右栏内容 / Right Column"
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    p = right_frame.add_paragraph()
    p.text = "• 要点一\n• 要点二\n• 要点三"
    p.font.name = "Arial"
    p.font.size = Pt(22)
    p.font.color.rgb = DARK_TEXT
    
    # Slide 4: Conclusion Slide
    print("  Creating conclusion slide layout...")
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    title = slide.shapes.title
    title.text = "结论 / Conclusion"
    title.text_frame.paragraphs[0].font.name = "Arial"
    title.text_frame.paragraphs[0].font.size = Pt(56)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    
    p = text_frame.paragraphs[0]
    p.text = "谢谢 / Thank You"
    p.font.name = "Arial"
    p.font.size = Pt(48)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    print(f"Saving master template to {output_path}...")
    prs.save(output_path)
    print("Master template creation complete!")

if __name__ == "__main__":
    # Transform the original PPT
    input_file = "original-strong-field-backup.pptx"
    output_file = "tech-business-template.pptx"
    
    if os.path.exists(input_file):
        apply_tech_business_style(input_file, output_file)
        print(f"\n✓ Transformed PPT saved as {output_file}")
    else:
        print(f"Warning: {input_file} not found, creating new template...")
        create_master_template(output_file)
        print(f"\n✓ New template created as {output_file}")
    
    # Also create a comprehensive master template
    master_template = "tech-business-master-template.pptx"
    create_master_template(master_template)
    print(f"\n✓ Master template created as {master_template}")
